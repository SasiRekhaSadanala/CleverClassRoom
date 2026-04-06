import os
import time
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_classic.chains import create_retrieval_chain
from langchain_classic.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from dotenv import load_dotenv
from pptx import Presentation

load_dotenv()

# We will use FAISS purely locally, nested by classroom
VECTOR_STORE_ROOT = "./faiss_store"
os.makedirs(VECTOR_STORE_ROOT, exist_ok=True)

def get_llm():
    return ChatGoogleGenerativeAI(model="models/gemini-2.0-flash", temperature=0)

def get_embeddings():
    return GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")

def load_pptx(file_path: str):
    """Custom loader for PPTX files."""
    prs = Presentation(file_path)
    text_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = f"--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text_content.append(slide_text)
    return "\n".join(text_content)

def process_and_index_pdf(file_path: str, classroom_id: int):
    """Loads a PDF or PPTX, chunks it, and indexes it into a classroom-specific FAISS store."""
    classroom_path = os.path.join(VECTOR_STORE_ROOT, f"class_{classroom_id}")
    
    if file_path.endswith('.pdf'):
        loader = PyPDFLoader(file_path)
        docs = loader.load()
    elif file_path.endswith('.pptx') or file_path.endswith('.ppt'):
        from langchain_core.documents import Document
        content = load_pptx(file_path)
        docs = [Document(page_content=content, metadata={"source": file_path})]
    else:
        return 0
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=4000, chunk_overlap=400)
    splits = text_splitter.split_documents(docs)
    
    embeddings = get_embeddings()
    vectorstore = None
    if os.path.exists(classroom_path):
        vectorstore = FAISS.load_local(classroom_path, embeddings, allow_dangerous_deserialization=True)
            
    batch_size = 20
    for i in range(0, len(splits), batch_size):
        batch = splits[i:i+batch_size]
        if vectorstore is None:
            vectorstore = FAISS.from_documents(batch, embeddings)
        else:
            vectorstore.add_documents(batch)
            
    if vectorstore:
        vectorstore.save_local(classroom_path)
        
    return len(splits)

def get_answer(question: str, classroom_id: int) -> str:
    """Answers a question using RAG from indexed documents for a specific classroom."""
    classroom_path = os.path.join(VECTOR_STORE_ROOT, f"class_{classroom_id}")
    
    if not os.path.exists(classroom_path):
        return "The professor hasn't uploaded any materials for this class yet."
        
    embeddings = get_embeddings()
    vectorstore = FAISS.load_local(classroom_path, embeddings, allow_dangerous_deserialization=True)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 5})
    
    llm = get_llm()
    system_prompt = (
        "You are an AI-powered Student Mentor for a specific classroom. Use the following pieces of retrieved context "
        "from the professor's uploaded materials to answer the student's question.\n\n"
        "STRICT RULE: ONLY use information from the provided context. If the answer is not in the context, "
        "politely say that the provided materials do not cover this topic and suggest they stick to the curriculum.\n\n"
        "Context:\n{context}"
    )
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("human", "{input}"),
    ])
    
    question_answer_chain = create_stuff_documents_chain(llm, prompt)
    rag_chain = create_retrieval_chain(retriever, question_answer_chain)
    
    try:
        response = rag_chain.invoke({"input": question})
        return response["answer"]
    except Exception as e:
        error_msg = str(e)
        if "429" in error_msg or "RESOURCE_EXHAUSTED" in error_msg:
            return "The AI mentor is currently at its capacity limit. Please wait about 60 seconds and try again."
        return f"I encountered an error while searching your materials: {error_msg}"
